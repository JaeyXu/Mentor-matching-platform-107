"""Evidence-first PowerPoint renderer for reviewed progress reports."""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

W, H = 13.333, 7.5
NAVY = RGBColor(15, 23, 42)
INDIGO = RGBColor(79, 70, 229)
TEAL = RGBColor(13, 148, 136)
SKY = RGBColor(14, 165, 233)
AMBER = RGBColor(245, 158, 11)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, value: str, left: float, top: float, width: float, height: float,
             size: int, color: RGBColor = NAVY, bold: bool = False,
             align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, left: float, top: float, width: float, height: float,
             fill: RGBColor, rounded: bool = False, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def new_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    return slide


def add_header(slide, title: str, subtitle: str) -> None:
    add_rect(slide, 0, 0, W, 0.16, INDIGO)
    add_text(slide, title, 0.7, 0.45, 9, 0.45, 25, NAVY, True)
    add_text(slide, subtitle, 0.72, 0.97, 11, 0.28, 10, MUTED)


def number(value: Any) -> int:
    try:
        return max(0, int(value or 0))
    except (TypeError, ValueError):
        return 0


def shorten(value: Any, limit: int = 34) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text if len(text) <= limit else text[:limit - 1] + "…"


def get_visual(payload: dict[str, Any]) -> dict[str, Any]:
    value = payload.get("visualData")
    return value if isinstance(value, dict) else {}


def metric(visual: dict[str, Any], key: str) -> int:
    metrics = visual.get("metrics")
    return number(metrics.get(key) if isinstance(metrics, dict) else 0)


def footer(slide, evidence: dict[str, int], page: int) -> None:
    source = " · ".join(
        f"{label}{count}"
        for label, count in (("计划", evidence.get("plan", 0)), ("活动", evidence.get("activity", 0)),
                             ("讨论", evidence.get("chat", 0)), ("其他", evidence.get("other", 0)))
        if count
    ) or "本期暂无可引用记录"
    add_text(slide, f"证据来源：{source}", 0.7, 7.05, 8.5, 0.24, 9, MUTED)
    add_text(slide, f"{page:02d}", 12.05, 7.03, 0.5, 0.24, 10, INDIGO, True, PP_ALIGN.RIGHT)


def add_cover(presentation: Presentation, title: str, visual: dict[str, Any], evidence: dict[str, int]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    period = visual.get("period") if isinstance(visual.get("period"), dict) else {}
    date_range = f"{period.get('start', '')} — {period.get('end', '')}".strip(" —")
    add_rect(slide, 0.8, 0.8, 1.3, 0.1, TEAL, True)
    add_text(slide, "RESEARCH WORKBENCH", 0.8, 1.2, 5.4, 0.3, 12, RGBColor(165, 180, 252), True)
    add_text(slide, title, 0.8, 1.82, 11, 0.9, 34, WHITE, True)
    add_text(slide, date_range or "基于本期已审核科研记录", 0.82, 2.85, 10.5, 0.35, 15, RGBColor(203, 213, 225))
    cards = (("可核验活动", metric(visual, "activity_events"), TEAL),
             ("完成计划", metric(visual, "completed_plans"), SKY),
             ("开放计划", metric(visual, "pending_plans"), AMBER))
    for index, (label, value, color) in enumerate(cards):
        x = 0.82 + index * 2.35
        add_rect(slide, x, 4.35, 2.05, 1.22, RGBColor(30, 41, 59), True)
        add_text(slide, str(value), x + 0.18, 4.58, 1.6, 0.38, 25, color, True)
        add_text(slide, label, x + 0.18, 5.02, 1.75, 0.25, 10, RGBColor(203, 213, 225))
    add_text(slide, "只呈现已记录、可复核的事实；未记录内容不作推断。", 0.82, 6.58, 10.5, 0.3, 11, RGBColor(148, 163, 184))
    footer(slide, evidence, 1)


def add_dashboard(presentation: Presentation, visual: dict[str, Any], evidence: dict[str, int], page: int) -> None:
    slide = new_slide(presentation)
    add_header(slide, "本期科研数据看板", "所有数字来自报告生成时锁定的可核验记录")
    cards = [
        ("科研活动", metric(visual, "activity_events"), "本期平台记录", INDIGO),
        ("完成计划", metric(visual, "completed_plans"), "已记录完成", TEAL),
        ("开放计划", metric(visual, "pending_plans"), "待办或进行中", AMBER),
        ("完成反馈", metric(visual, "completion_feedbacks"), "已记录成果", SKY),
        ("论文阅读", metric(visual, "read_papers"), "已记录阅读", INDIGO),
        ("导师匹配", metric(visual, "matched_mentors"), "已审核匹配", TEAL),
    ]
    for index, (label, value, note, color) in enumerate(cards):
        col, row = index % 3, index // 3
        x, y = 0.75 + col * 4.15, 1.52 + row * 1.62
        add_rect(slide, x, y, 3.75, 1.28, WHITE, True, LINE)
        add_rect(slide, x, y, 0.08, 1.28, color)
        add_text(slide, label, x + 0.26, y + 0.2, 2.3, 0.25, 12, SLATE, True)
        add_text(slide, str(value), x + 0.26, y + 0.52, 1.0, 0.4, 25, NAVY, True)
        add_text(slide, note, x + 1.25, y + 0.63, 2.1, 0.24, 10, MUTED)
    done, open_count = metric(visual, "completed_plans"), metric(visual, "pending_plans")
    ratio = done / max(1, done + open_count)
    add_text(slide, "计划完成比例", 0.78, 5.18, 2.2, 0.25, 13, NAVY, True)
    add_rect(slide, 0.78, 5.6, 8.55, 0.28, LINE, True)
    if ratio:
        add_rect(slide, 0.78, 5.6, 8.55 * ratio, 0.28, TEAL, True)
    add_text(slide, f"{done} 已完成 / {open_count} 开放（{round(ratio * 100)}%）", 0.78, 5.98, 4.8, 0.25, 11, MUTED)
    add_text(slide, "阅读建议", 10.0, 5.18, 1.5, 0.25, 13, NAVY, True)
    add_text(slide, "先看本期变化，再看计划状态，最后落实下一步交付物。", 10.0, 5.58, 2.3, 0.7, 12, SLATE)
    footer(slide, evidence, page)


def add_trend(presentation: Presentation, visual: dict[str, Any], evidence: dict[str, int], page: int) -> None:
    slide = new_slide(presentation)
    add_header(slide, "科研活动变化", "按同类报告的历史快照展示，不是模型推断")
    history = visual.get("history") if isinstance(visual.get("history"), list) else []
    rows = [item for item in history if isinstance(item, dict)][-7:]
    if not rows:
        add_text(slide, "尚无历史报告快照；生成下一期报告后将形成可比较趋势。", 0.85, 2.7, 10, 0.4, 17, MUTED)
        footer(slide, evidence, page)
        return
    maximum = max(1, *(number(item.get("activity_events")) for item in rows), *(number(item.get("completed_plans")) for item in rows))
    left, top, width, height = 0.95, 1.7, 10.9, 4.0
    for tick in range(5):
        y = top + height - tick * height / 4
        add_rect(slide, left, y, width, 0.01, LINE)
        add_text(slide, str(round(maximum * tick / 4)), 0.35, y - 0.1, 0.45, 0.2, 9, MUTED, align=PP_ALIGN.RIGHT)
    step = width / max(1, len(rows))
    for index, item in enumerate(rows):
        x = left + index * step + step * 0.22
        activity, complete = number(item.get("activity_events")), number(item.get("completed_plans"))
        add_rect(slide, x, top + height - height * activity / maximum, step * 0.25, max(0.03, height * activity / maximum), INDIGO, True)
        add_rect(slide, x + step * 0.31, top + height - height * complete / maximum, step * 0.25, max(0.03, height * complete / maximum), TEAL, True)
        add_text(slide, shorten(item.get("label"), 8), x - 0.06, top + height + 0.15, step * 0.7, 0.22, 9, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.0, 6.28, 0.17, 0.17, INDIGO, True)
    add_text(slide, "科研活动记录", 1.24, 6.22, 1.7, 0.25, 10, SLATE)
    add_rect(slide, 3.1, 6.28, 0.17, 0.17, TEAL, True)
    add_text(slide, "完成计划", 3.34, 6.22, 1.4, 0.25, 10, SLATE)
    if len(rows) == 1:
        add_text(slide, "当前只有 1 期快照，图表用于建立后续比较基线。", 7.0, 6.22, 4.6, 0.25, 10, MUTED, align=PP_ALIGN.RIGHT)
    footer(slide, evidence, page)


def set_cell(cell, value: str, fill: RGBColor, color: RGBColor, bold: bool = False) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.text = ""
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(10)
    run.font.bold = bold
    run.font.color.rgb = color


def add_plan_table(presentation: Presentation, visual: dict[str, Any], evidence: dict[str, int], page: int) -> None:
    slide = new_slide(presentation)
    add_header(slide, "计划状态与时间预算", "以报告期结束时的计划记录为准")
    plans = visual.get("plans") if isinstance(visual.get("plans"), list) else []
    rows = [item for item in plans if isinstance(item, dict)][:5]
    if not rows:
        add_text(slide, "本期没有可展示的计划记录。", 0.85, 2.7, 8, 0.4, 17, MUTED)
        footer(slide, evidence, page)
        return
    table = slide.shapes.add_table(len(rows) + 1, 5, Inches(0.75), Inches(1.65), Inches(11.85), Inches(4.6)).table
    for index, width in enumerate((4.45, 1.25, 1.35, 2.35, 2.45)):
        table.columns[index].width = Inches(width)
    for index, value in enumerate(("计划", "状态", "优先级", "截止时间", "投入 / 实际")):
        set_cell(table.cell(0, index), value, NAVY, WHITE, True)
    status = {"todo": "待办", "doing": "进行中", "done": "已完成", "cancelled": "已取消"}
    priority = {"high": "高", "medium": "中", "low": "低"}
    for row_index, item in enumerate(rows, start=1):
        fill = WHITE if row_index % 2 else RGBColor(241, 245, 249)
        set_cell(table.cell(row_index, 0), shorten(item.get("title"), 44), fill, NAVY, True)
        set_cell(table.cell(row_index, 1), status.get(str(item.get("status")), "未设置"), fill, TEAL if item.get("status") == "done" else SLATE)
        set_cell(table.cell(row_index, 2), priority.get(str(item.get("priority")), "中"), fill, SLATE)
        set_cell(table.cell(row_index, 3), shorten(str(item.get("due_at") or "未设置").replace("T", " "), 18), fill, SLATE)
        set_cell(table.cell(row_index, 4), f"{number(item.get('estimated_minutes'))} / {number(item.get('actual_minutes'))} 分钟", fill, SLATE)
    footer(slide, evidence, page)


def action_items(markdown: str) -> list[str]:
    capture, values = False, []
    for raw in markdown.splitlines():
        line = raw.strip()
        if line.startswith("## 七、"):
            capture = True
            continue
        if capture and line.startswith("## "):
            break
        if capture and re.match(r"^\d+\.\s", line):
            line = re.sub(r"^\d+\.\s+", "", line)
            line = re.sub(r"\*\*", "", line)
            values.append(shorten(re.sub(r"^\[[A-Z]+\]\s*", "", line), 86))
    return values[:3]


def add_actions(presentation: Presentation, markdown: str, visual: dict[str, Any], evidence: dict[str, int], page: int) -> None:
    slide = new_slide(presentation)
    add_header(slide, "下一步：可验收行动", "只展示报告中已有的后续行动，不自动声称任务已完成")
    values = action_items(markdown) or ["当前报告未形成可展示的后续行动，请补充一项有交付物和验收标准的计划。"]
    for index, item in enumerate(values):
        y = 1.65 + index * 1.42
        color = (INDIGO, TEAL, SKY)[index % 3]
        add_rect(slide, 0.85, y, 11.55, 1.05, WHITE, True, LINE)
        add_rect(slide, 0.85, y, 0.72, 1.05, color, True)
        add_text(slide, str(index + 1), 1.02, y + 0.28, 0.35, 0.35, 18, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, item, 1.83, y + 0.2, 9.95, 0.62, 16, NAVY, True)
    add_text(slide, f"本期行动基线：{metric(visual, 'pending_plans')} 项开放计划，{metric(visual, 'completed_plans')} 项完成计划。", 0.87, 6.13, 7.2, 0.3, 11, MUTED)
    add_text(slide, "汇报建议：优先说明交付物、验收结果和下一步问题，而不是复述计划标题。", 0.87, 6.48, 10.9, 0.28, 11, SLATE)
    footer(slide, evidence, page)


def main() -> int:
    if len(sys.argv) != 2:
        print("output path required", file=sys.stderr)
        return 2
    payload = json.loads(sys.stdin.read().lstrip("\ufeff"))
    output = Path(sys.argv[1])
    title = str(payload.get("title") or "科研进展报告")
    slide_count = max(3, min(20, int(payload.get("slideCount") or 5)))
    visual = get_visual(payload)
    refs = [str(item) for item in payload.get("evidenceRefs") or [] if str(item)]
    raw_evidence = visual.get("evidence") if isinstance(visual.get("evidence"), dict) else {}
    evidence = {key: number(raw_evidence.get(key)) for key in ("plan", "activity", "chat", "other")}
    if not sum(evidence.values()) and refs:
        evidence["other"] = len(refs)

    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    add_cover(presentation, title, visual, evidence)
    builders = [
        lambda page: add_dashboard(presentation, visual, evidence, page),
        lambda page: add_trend(presentation, visual, evidence, page),
        lambda page: add_plan_table(presentation, visual, evidence, page),
        lambda page: add_actions(presentation, str(payload.get("markdown") or ""), visual, evidence, page),
    ]
    for page, builder in enumerate(builders[:max(0, slide_count - 1)], start=2):
        builder(page)
    presentation.save(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
